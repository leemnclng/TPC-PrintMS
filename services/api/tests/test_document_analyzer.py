from io import BytesIO

from docx import Document
from fastapi import FastAPI
from fastapi.testclient import TestClient
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter
import pymupdf
import pytest
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from app.core.config import settings
from app.db.base import Base
from app.db.models import ProductPrintType
from app.db.session import get_db
from app.modules.document_analyzer.api import router
from app.modules.document_analyzer.models.document_analysis import DocumentAnalysis
from app.modules.document_analyzer.models.enums import DocumentFileType, Orientation, PaperSize
from app.modules.document_analyzer.pricing.calculator import calculate_price
from app.modules.document_analyzer.services.analysis_service import AnalysisService
from app.routers import inventory, job_orders, products, services, variants


def test_supported_document_formats_produce_normalized_analysis() -> None:
    service = AnalysisService()
    fixtures = _fixtures()

    image = service.analyze("colored-a4.png", fixtures["image"], "image/png")
    assert image.file_type.value == "image"
    assert image.paper_size.value == "A4"
    assert image.color_pages == 1
    assert image.ocr_required is True

    pdf = service.analyze("blank.pdf", fixtures["pdf"], "application/pdf")
    assert pdf.file_type.value == "pdf"
    assert pdf.page_count == 1
    assert pdf.paper_size.value == "A4"

    colored_pdf = service.analyze("colored.pdf", fixtures["colored_pdf"], "application/pdf")
    assert colored_pdf.color_pages == 1
    assert colored_pdf.estimated_color_coverage_percent > 95
    assert colored_pdf.estimated_ink_coverage_percent > 95

    docx = service.analyze("sample.docx", fixtures["docx"])
    assert docx.word_count == 5
    assert docx.table_count == 1

    xlsx = service.analyze("sample.xlsx", fixtures["xlsx"])
    assert xlsx.page_count == 1
    assert xlsx.word_count == 2

    pptx = service.analyze("sample.pptx", fixtures["pptx"])
    assert pptx.page_count == 1
    assert pptx.word_count == 2


@pytest.mark.parametrize("paper_size", [PaperSize.a4, PaperSize.letter, PaperSize.legal])
def test_bw_product_rate_already_includes_paper_and_ink_for_every_size(paper_size: PaperSize) -> None:
    analysis = DocumentAnalysis(
        filename="four-pages.pdf",
        file_type=DocumentFileType.pdf,
        mime_type="application/pdf",
        file_size_bytes=1024,
        page_count=4,
        paper_size=paper_size,
        orientation=Orientation.portrait,
        character_count=0,
        word_count=0,
        ocr_required=False,
        image_count=0,
        contains_images=False,
        image_coverage_percent=0,
        estimated_color_coverage_percent=90,
        estimated_ink_coverage_percent=80,
        table_count=0,
        graphic_count=0,
        color_pages=4,
        bw_pages=0,
        duplex_compatible=True,
        estimated_print_time_seconds=4,
        confidence=1,
    )

    pricing = calculate_price(
        analysis,
        {
            ProductPrintType.black_and_white: (3, "paperSize"),
            ProductPrintType.colored: (10, "paperSize"),
        },
        ProductPrintType.black_and_white,
    )

    assert pricing.base_subtotal == 12
    assert pricing.suggested_price == 12
    assert pricing.adjustments == []


def test_analyzer_api_and_owner_pricing_rules(tmp_path) -> None:
    engine = create_engine(
        f"sqlite:///{tmp_path / 'analyzer.db'}",
        connect_args={"check_same_thread": False},
    )
    test_session = sessionmaker(autocommit=False, autoflush=False, bind=engine)
    Base.metadata.create_all(engine)

    def override_db():
        db = test_session()
        try:
            yield db
        finally:
            db.close()

    app = FastAPI()
    app.dependency_overrides[get_db] = override_db
    app.include_router(router)
    app.include_router(inventory.router)
    client = TestClient(app)
    headers = {"X-Print-MS-Token": settings.token}
    image = _fixtures()["image"]

    # Paper sizes are tied to real inventory stock — see decisions.md "Tie
    # Document Pricing to Real Paper Stock" — so nothing prices until an
    # item is tagged, and a freshly-tagged item's rule starts at ₱0.
    client.post(
        "/inventory-items",
        headers=headers,
        json={
            "name": "A4 Bond Paper",
            "category": "Paper",
            "unit": "sheet",
            "openingQuantity": 100,
            "reorderLevel": 10,
            "paperSize": "A4",
            "isActive": True,
        },
    )

    rules_response = client.get("/document-analyzer/pricing-rules", headers=headers)
    assert rules_response.status_code == 200
    rules = rules_response.json()
    assert len(rules) == 8
    assert {rule["pricingScope"] for rule in rules} == {"printing", "photocopy"}
    a4_color = next(rule for rule in rules if rule["paperSize"] == "A4" and rule["printType"] == "colored" and rule["pricingScope"] == "printing")
    assert a4_color["pricePerPage"] == 0
    seed_response = client.put(
        "/document-analyzer/pricing-rules",
        headers=headers,
        json={"rules": [{"id": a4_color["id"], "pricePerPage": 5, "isActive": True}]},
    )
    assert seed_response.status_code == 200

    response = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        files={"file": ("colored-a4.png", image, "image/png")},
    )
    assert response.status_code == 200
    result = response.json()
    assert result["analysis"]["paperSize"] == "A4"
    assert result["analysis"]["colorPages"] == 1
    assert result["pricing"]["suggestedPrice"] == 5
    assert result["pricing"]["currency"] == "PHP"

    update_response = client.put(
        "/document-analyzer/pricing-rules",
        headers=headers,
        json={"rules": [{"id": a4_color["id"], "pricePerPage": 7.5, "isActive": True}]},
    )
    assert update_response.status_code == 200

    repriced = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        files={"file": ("colored-a4.png", image, "image/png")},
    ).json()
    assert repriced["pricing"]["suggestedPrice"] == 7.5

    unsupported = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        files={"file": ("notes.txt", b"hello", "text/plain")},
    )
    assert unsupported.status_code == 415

    invalid_rate = client.put(
        "/document-analyzer/pricing-rules",
        headers=headers,
        json={"rules": [{"id": a4_color["id"], "pricePerPage": -1, "isActive": True}]},
    )
    assert invalid_rate.status_code == 422


def test_analyze_prefers_product_override_then_falls_back_to_global_rate(tmp_path) -> None:
    engine = create_engine(
        f"sqlite:///{tmp_path / 'analyzer_product.db'}",
        connect_args={"check_same_thread": False},
    )
    test_session = sessionmaker(autocommit=False, autoflush=False, bind=engine)
    Base.metadata.create_all(engine)

    def override_db():
        db = test_session()
        try:
            yield db
        finally:
            db.close()

    app = FastAPI()
    app.dependency_overrides[get_db] = override_db
    app.include_router(router)
    app.include_router(services.router)
    app.include_router(products.router)
    app.include_router(inventory.router)
    app.include_router(variants.router)
    client = TestClient(app)
    headers = {"X-Print-MS-Token": settings.token}
    image = _fixtures()["image"]

    service = client.post(
        "/services",
        headers=headers,
        json={"name": "Photo printing", "category": "printing", "description": None, "isActive": True},
    ).json()
    material = client.post(
        "/inventory-items",
        headers=headers,
        json={
            "name": "Photo paper",
            "category": "Paper",
            "unit": "sheet",
            "openingQuantity": 50,
            "reorderLevel": 5,
            "paperSize": "A4",
            "isActive": True,
        },
    ).json()

    rules = client.get("/document-analyzer/pricing-rules", headers=headers).json()
    a4_color_rule = next(rule for rule in rules if rule["paperSize"] == "A4" and rule["printType"] == "colored" and rule["pricingScope"] == "printing")
    a4_bw_rule = next(rule for rule in rules if rule["paperSize"] == "A4" and rule["printType"] == "black_and_white" and rule["pricingScope"] == "printing")
    a4_semi_rule = next(rule for rule in rules if rule["paperSize"] == "A4" and rule["printType"] == "semi_colored" and rule["pricingScope"] == "printing")
    seed_response = client.put(
        "/document-analyzer/pricing-rules",
        headers=headers,
        json={
            "rules": [
                {"id": a4_color_rule["id"], "pricePerPage": 5, "isActive": True},
                {"id": a4_bw_rule["id"], "pricePerPage": 2, "isActive": True},
                {"id": a4_semi_rule["id"], "pricePerPage": 8, "isActive": True},
            ]
        },
    )
    assert seed_response.status_code == 200

    variant = client.post(
        "/variants",
        headers=headers,
        json={"label": "Back-to-back", "description": None, "isActive": True},
    ).json()

    product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Premium photo print",
            "printType": "colored",
            "isActive": True,
            "variants": [{"variantId": variant["id"], "priceAdjustment": 2}],
            "materialAssignments": [{"inventoryItemId": material["id"]}],
            "documentRates": [{"pricingRuleId": a4_color_rule["id"], "pricePerPage": 15}],
        },
    ).json()
    # The product reference follows its assigned A4 material and prefers the
    # product-specific override over the global rate.
    assert product["pricePerPage"] == 15
    assert product["documentRates"] == [
        {
            "id": product["documentRates"][0]["id"],
            "pricingRuleId": a4_color_rule["id"],
            "pricePerPage": 15,
            "paperSize": "A4",
            "printType": "colored",
            "pricingScope": "printing",
        }
    ]

    priced = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": product["id"]},
        files={"file": ("colored-a4.png", image, "image/png")},
    ).json()
    assert priced["pricingContext"]["productId"] == product["id"]
    assert priced["pricingContext"]["productName"] == "Premium photo print"
    breakdown = priced["pricing"]["breakdown"][0]
    assert breakdown["ratePerPage"] == 15
    assert breakdown["rateSource"] == "product"
    assert priced["analysis"]["estimatedInkCoveragePercent"] == 88.2
    assert priced["pricing"]["baseSubtotal"] == 15
    assert priced["pricing"]["adjustments"][0]["kind"] == "inkCoverage"
    assert priced["pricing"]["suggestedPrice"] == 28.23

    sibling_product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Same paper, different product price",
            "printType": "colored",
            "isActive": True,
            "variants": [],
            "materialAssignments": [{"inventoryItemId": material["id"]}],
            "documentRates": [{"pricingRuleId": a4_color_rule["id"], "pricePerPage": 22}],
        },
    ).json()
    sibling_priced = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": sibling_product["id"]},
        files={"file": ("same-colored-a4.png", image, "image/png")},
    ).json()
    assert sibling_priced["pricing"]["baseSubtotal"] == 22
    assert client.get(f"/products/{product['id']}", headers=headers).json()["pricePerPage"] == 15

    with_variant = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": product["id"], "variant_id": variant["id"]},
        files={"file": ("colored-a4.png", image, "image/png")},
    ).json()
    assert with_variant["pricingContext"]["variantName"] == "Back-to-back"
    assert with_variant["pricing"]["adjustments"][-1]["kind"] == "variant"
    assert with_variant["pricing"]["suggestedPrice"] == 30.23

    bw_product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Standard document",
            "printType": "black_and_white",
            "isActive": True,
            "variants": [],
            "materialAssignments": [{"inventoryItemId": material["id"]}],
            "documentRates": [],
        },
    ).json()
    bw_with_color = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": bw_product["id"]},
        files={"file": ("colored-a4.png", image, "image/png")},
    ).json()
    assert bw_with_color["pricing"]["baseSubtotal"] == 2
    assert bw_with_color["pricing"]["adjustments"] == []
    assert bw_with_color["pricing"]["suggestedPrice"] == 2

    semi_product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Semi-colored document",
            "printType": "semi_colored",
            "isActive": True,
            "variants": [],
            "materialAssignments": [{"inventoryItemId": material["id"]}],
            "documentRates": [],
        },
    ).json()
    semi_priced = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": semi_product["id"]},
        files={"file": ("colored-a4.png", image, "image/png")},
    ).json()
    assert semi_priced["pricing"]["baseSubtotal"] == 8
    assert semi_priced["pricing"]["adjustments"][0]["kind"] == "inkCoverage"
    assert semi_priced["pricing"]["suggestedPrice"] == 15.06

    # A selected product defines how every page is printed. The analyzer's
    # detected page color is still reported, but it must not switch away
    # from that product's configured print type or assigned paper price.
    bw_pdf_priced_as_color_product = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": product["id"]},
        files={"file": ("blank.pdf", _fixtures()["pdf"], "application/pdf")},
    ).json()
    assert bw_pdf_priced_as_color_product["analysis"]["bwPages"] == 1
    assert bw_pdf_priced_as_color_product["pricing"]["suggestedPrice"] == 15
    assert bw_pdf_priced_as_color_product["pricing"]["breakdown"][0]["printType"] == "colored"

    generic = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        files={"file": ("colored-a4.png", image, "image/png")},
    ).json()
    assert generic.get("pricingContext") is None
    assert generic["pricing"]["breakdown"][0]["rateSource"] == "paperSize"
    assert generic["pricing"]["suggestedPrice"] == 5

    missing_product = client.post(
        "/document-analyzer/analyze",
        headers=headers,
        data={"product_id": "missing"},
        files={"file": ("colored-a4.png", image, "image/png")},
    )
    assert missing_product.status_code == 404
    duplicate_override = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Duplicate override",
            "printType": "colored",
            "isActive": True,
            "variants": [],
            "materialAssignments": [{"inventoryItemId": material["id"]}],
            "documentRates": [
                {"pricingRuleId": a4_color_rule["id"], "pricePerPage": 12},
                {"pricingRuleId": a4_color_rule["id"], "pricePerPage": 13},
            ],
        },
    )
    assert duplicate_override.status_code == 409


def test_scan_pricing_tiers_by_page_count_and_product_override(tmp_path) -> None:
    engine = create_engine(
        f"sqlite:///{tmp_path / 'scan_pricing.db'}",
        connect_args={"check_same_thread": False},
    )
    test_session = sessionmaker(autocommit=False, autoflush=False, bind=engine)
    Base.metadata.create_all(engine)

    def override_db():
        db = test_session()
        try:
            yield db
        finally:
            db.close()

    app = FastAPI()
    app.dependency_overrides[get_db] = override_db
    app.include_router(router)
    app.include_router(services.router)
    app.include_router(products.router)
    app.include_router(job_orders.router)
    client = TestClient(app)
    headers = {"X-Print-MS-Token": settings.token}

    def scan_page() -> bytes:
        buffer = BytesIO()
        Image.new("RGB", (794, 1123), "white").save(buffer, format="PNG")
        return buffer.getvalue()

    service = client.post(
        "/services",
        headers=headers,
        json={"name": "Scan bureau", "category": "photocopy", "description": None, "isActive": True},
    ).json()
    product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Document scan",
            "operationKind": "scan",
            "printType": "black_and_white",
            "isActive": True,
            "variants": [],
            "materialAssignments": [],
            "documentRates": [],
        },
    ).json()
    assert product["standalonePricePerPage"] is None

    # No tiers configured yet, and no product override: nothing to resolve.
    assert client.get("/document-analyzer/scan-pricing-tiers", headers=headers).json() == []
    unpriced_job = client.post(
        "/job-orders/from-scan",
        headers=headers,
        json={"name": "Unpriced scan job", "serviceId": service["id"], "productId": product["id"]},
    )
    assert unpriced_job.status_code == 422

    tier_1_to_5 = client.post(
        "/document-analyzer/scan-pricing-tiers",
        headers=headers,
        json={"minPages": 1, "maxPages": 5, "pricePerPage": 10, "isActive": True},
    ).json()
    assert tier_1_to_5["minPages"] == 1
    assert tier_1_to_5["maxPages"] == 5

    # An overlapping range is rejected.
    overlap_response = client.post(
        "/document-analyzer/scan-pricing-tiers",
        headers=headers,
        json={"minPages": 4, "maxPages": 8, "pricePerPage": 8, "isActive": True},
    )
    assert overlap_response.status_code == 409

    # An open-ended top tier ("6 and up") is fine once it starts after tier 1.
    tier_6_and_up = client.post(
        "/document-analyzer/scan-pricing-tiers",
        headers=headers,
        json={"minPages": 6, "maxPages": None, "pricePerPage": 8, "isActive": True},
    ).json()
    assert tier_6_and_up["maxPages"] is None

    tiers = client.get("/document-analyzer/scan-pricing-tiers", headers=headers).json()
    assert {tier["id"] for tier in tiers} == {tier_1_to_5["id"], tier_6_and_up["id"]}

    # A job can be created now that pricing exists somewhere, even before the
    # real page count (and therefore the exact rate) is known.
    create_response = client.post(
        "/job-orders/from-scan",
        headers=headers,
        json={"name": "Tiered scan job", "serviceId": service["id"], "productId": product["id"]},
    )
    assert create_response.status_code == 201
    job_response = create_response.json()

    # 3 pages falls in the 1-5 tier: ₱10 × 3 = ₱30.
    small_scan = client.post(
        f"/job-orders/{job_response['id']}/scan-output",
        headers=headers,
        files=[("files", ("page.png", scan_page(), "image/png")) for _ in range(3)],
    ).json()
    assert small_scan["items"][0]["pagesPerCopy"] == 3
    assert small_scan["items"][0]["unitPrice"] == 10
    assert small_scan["total"] == 30

    # A re-scan landing in the open-ended tier switches rates: 8 pages × ₱8 = ₱64.
    requeue = client.post(
        f"/job-orders/{job_response['id']}/transitions",
        headers=headers,
        json={"toStatus": "queued"},
    )
    assert requeue.status_code == 200
    big_scan = client.post(
        f"/job-orders/{job_response['id']}/scan-output",
        headers=headers,
        files=[("files", ("page.png", scan_page(), "image/png")) for _ in range(8)],
    ).json()
    assert big_scan["items"][0]["pagesPerCopy"] == 8
    assert big_scan["items"][0]["unitPrice"] == 8
    assert big_scan["total"] == 64

    # A product's own price always wins over the tier table, regardless of
    # how many pages it ends up being.
    overridden_product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Flat-rate scan",
            "operationKind": "scan",
            "printType": "black_and_white",
            "standalonePricePerPage": 20,
            "isActive": True,
            "variants": [],
            "materialAssignments": [],
            "documentRates": [],
        },
    ).json()
    override_create_response = client.post(
        "/job-orders/from-scan",
        headers=headers,
        json={"name": "Flat-rate scan job", "serviceId": service["id"], "productId": overridden_product["id"]},
    )
    assert override_create_response.status_code == 201
    override_job = override_create_response.json()
    override_scan = client.post(
        f"/job-orders/{override_job['id']}/scan-output",
        headers=headers,
        files=[("files", ("page.png", scan_page(), "image/png")) for _ in range(2)],
    ).json()
    assert override_scan["items"][0]["unitPrice"] == 20
    assert override_scan["total"] == 40

    # A page count that falls in a gap between configured tiers has nothing
    # to resolve, so the scan is rejected rather than priced at zero.
    gapped_product = client.post(
        "/products",
        headers=headers,
        json={
            "serviceId": service["id"],
            "name": "Gapped scan",
            "operationKind": "scan",
            "printType": "black_and_white",
            "isActive": True,
            "variants": [],
            "materialAssignments": [],
            "documentRates": [],
        },
    ).json()
    gapped_create_response = client.post(
        "/job-orders/from-scan",
        headers=headers,
        json={"name": "Gapped scan job", "serviceId": service["id"], "productId": gapped_product["id"]},
    )
    assert gapped_create_response.status_code == 201
    gapped_job = gapped_create_response.json()
    client.put(
        f"/document-analyzer/scan-pricing-tiers/{tier_1_to_5['id']}",
        headers=headers,
        json={"minPages": 1, "maxPages": 2, "pricePerPage": 10, "isActive": True},
    )
    gapped_scan = client.post(
        f"/job-orders/{gapped_job['id']}/scan-output",
        headers=headers,
        files=[("files", ("page.png", scan_page(), "image/png")) for _ in range(4)],
    )
    assert gapped_scan.status_code == 422

    # Tiers can be removed entirely.
    delete_response = client.delete(f"/document-analyzer/scan-pricing-tiers/{tier_6_and_up['id']}", headers=headers)
    assert delete_response.status_code == 204
    remaining_tiers = client.get("/document-analyzer/scan-pricing-tiers", headers=headers).json()
    assert {tier["id"] for tier in remaining_tiers} == {tier_1_to_5["id"]}


def _fixtures() -> dict[str, bytes]:
    image = Image.new("RGB", (248, 351), (190, 30, 30))
    image_buffer = BytesIO()
    image.save(image_buffer, "PNG", dpi=(30, 30))

    pdf_writer = PdfWriter()
    pdf_writer.add_blank_page(width=595.28, height=841.89)
    pdf_buffer = BytesIO()
    pdf_writer.write(pdf_buffer)

    colored_pdf_document = pymupdf.open()
    colored_pdf_page = colored_pdf_document.new_page(width=595.28, height=841.89)
    colored_pdf_page.draw_rect(colored_pdf_page.rect, color=(1, 0, 0), fill=(1, 0, 0))
    colored_pdf = colored_pdf_document.tobytes()
    colored_pdf_document.close()

    document = Document()
    document.add_paragraph("Hello world from Word")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Table"
    docx_buffer = BytesIO()
    document.save(docx_buffer)

    workbook = Workbook()
    workbook.active["A1"] = "Hello workbook"
    xlsx_buffer = BytesIO()
    workbook.save(xlsx_buffer)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Hello slides"
    pptx_buffer = BytesIO()
    presentation.save(pptx_buffer)

    return {
        "image": image_buffer.getvalue(),
        "pdf": pdf_buffer.getvalue(),
        "colored_pdf": colored_pdf,
        "docx": docx_buffer.getvalue(),
        "xlsx": xlsx_buffer.getvalue(),
        "pptx": pptx_buffer.getvalue(),
    }
