from __future__ import annotations

from io import BytesIO
from statistics import mean

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models.document_analysis import DocumentAnalysis
from ..models.enums import DocumentFileType
from ..utils.color_analysis import ImageColorMetrics, analyze_image_color
from ..utils.image_processing import open_image
from ..utils.page_geometry import classify_orientation, classify_paper_size, emu_to_mm
from .base import DocumentAnalyzer, InvalidDocumentError, estimate_print_time


class PptxAnalyzer(DocumentAnalyzer):
    file_type = DocumentFileType.pptx

    def analyze(self, filename: str, data: bytes, mime_type: str) -> DocumentAnalysis:
        try:
            presentation = Presentation(BytesIO(data))
        except Exception as error:
            raise InvalidDocumentError("The selected PowerPoint file is damaged or unsupported.") from error
        if len(presentation.slides) == 0:
            raise InvalidDocumentError("The selected presentation has no slides.")

        text_parts: list[str] = []
        image_count = 0
        graphic_count = 0
        table_count = 0
        color_pages = 0
        image_metrics: list[ImageColorMetrics] = []
        for slide in presentation.slides:
            slide_colored = False
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text_parts.append(shape.text)
                    slide_colored = slide_colored or self._text_is_colored(shape)
                if getattr(shape, "has_table", False):
                    table_count += 1
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    try:
                        metric = analyze_image_color(open_image(shape.image.blob))
                        image_metrics.append(metric)
                        slide_colored = slide_colored or metric.is_colored
                    except Exception:
                        pass
                elif shape.shape_type not in {MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.GROUP}:
                    graphic_count += 1
                    slide_colored = slide_colored or self._shape_fill_is_colored(shape)
            color_pages += int(slide_colored)

        text = "\n".join(text_parts).strip()
        page_count = len(presentation.slides)
        bw_pages = page_count - color_pages
        width_mm = emu_to_mm(presentation.slide_width)
        height_mm = emu_to_mm(presentation.slide_height)
        coverage = round(mean(metric.coverage_percent for metric in image_metrics), 1) if image_metrics else 0.0
        ink = round(mean(metric.ink_coverage_percent for metric in image_metrics), 1) if image_metrics else 0.0
        return DocumentAnalysis(
            filename=filename,
            file_type=self.file_type,
            mime_type=mime_type or "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            file_size_bytes=len(data),
            page_count=page_count,
            paper_size=classify_paper_size(width_mm, height_mm),
            orientation=classify_orientation(width_mm, height_mm),
            width_mm=width_mm,
            height_mm=height_mm,
            dpi=None,
            character_count=len(text),
            word_count=len(text.split()),
            ocr_required=False,
            image_count=image_count,
            contains_images=image_count > 0,
            image_coverage_percent=coverage,
            estimated_ink_coverage_percent=ink,
            table_count=table_count,
            graphic_count=graphic_count,
            color_pages=color_pages,
            bw_pages=bw_pages,
            duplex_compatible=page_count > 1,
            estimated_print_time_seconds=estimate_print_time(color_pages, bw_pages),
            confidence=0.88,
            warnings=["PowerPoint pages represent slides; printer handout layouts are not included."],
        )

    @staticmethod
    def _text_is_colored(shape) -> bool:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    rgb = run.font.color.rgb
                    if rgb and max(rgb) - min(rgb) > 14:
                        return True
                except Exception:
                    continue
        return False

    @staticmethod
    def _shape_fill_is_colored(shape) -> bool:
        try:
            rgb = shape.fill.fore_color.rgb
            return bool(rgb and max(rgb) - min(rgb) > 14)
        except Exception:
            return False
